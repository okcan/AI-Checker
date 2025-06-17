import os
import spacy
import textstat
import argparse
import pdfplumber
from docx import Document
from pptx import Presentation

nlp = spacy.load("en_core_web_sm")

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".txt":
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

    elif ext == ".docx":
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == ".pdf":
        text = ""
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                content = page.extract_text()
                if content:
                    text += content + "\n"
        return text

    elif ext == ".pptx":
        prs = Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    else:
        raise ValueError(f"Unsupported file type: {ext}")

def ai_likeness_score(sentence: str) -> float:
    doc = nlp(sentence)
    length_score = min(len(sentence) / 200, 1)
    conj_score = sum(1 for tok in doc if tok.dep_ == "cc") / len(doc)
    words = [tok.text.lower() for tok in doc if tok.is_alpha]
    diversity = len(set(words)) / (len(words) + 1e-5)
    diversity_score = 1 - min(diversity, 1)
    try:
        read_score = 1 - min(textstat.flesch_reading_ease(sentence) / 100, 1)
    except:
        read_score = 0.5
    return (length_score + conj_score + diversity_score + read_score) / 4

def classify_score(score):
    if score > 0.65:
        return "⚠️ Likely AI-written"
    elif score > 0.45:
        return "🟡 Possibly AI-supported"
    else:
        return "✅ Likely Human"

def analyze_text(text: str):
    doc = nlp(text)
    sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
    results = []
    total_score = 0
    ai_like, mixed, human = 0, 0, 0

    for sent in sentences:
        score = ai_likeness_score(sent)
        total_score += score
        label = classify_score(score)
        results.append((sent, score, label))
        if "AI-written" in label:
            ai_like += 1
        elif "Possibly" in label:
            mixed += 1
        else:
            human += 1

    avg_score = total_score / len(sentences) if sentences else 0
    ai_percent = round(avg_score * 100, 2)

    # 🔍 Detaylı sonuçlar
    print(f"\n🧠 Estimated AI-style writing: {ai_percent}% of the text.\n")

    for sent, score, label in results:
        print(f"{label} ({round(score*100)}%) → {sent}")

    # 📊 Rapor
    print("\n--- Summary Report ---")
    print(f"Total sentences analyzed: {len(sentences)}")
    print(f"Average AI-likeness score : {ai_percent}%")
    print(f"✅ Human-like sentences   : {human}")
    print(f"🟡 Possibly AI-supported  : {mixed}")
    print(f"⚠️ AI-like sentences      : {ai_like}")

    # 🔍 Genel karar
    if ai_like / len(sentences) > 0.5:
        decision = "🧠 Final verdict: This text is likely AI-generated."
    elif (ai_like + mixed) / len(sentences) > 0.5:
        decision = "🧠 Final verdict: This text is possibly AI-assisted."
    else:
        decision = "🧠 Final verdict: This text is likely written by a human."

    print(decision)

def main():
    parser = argparse.ArgumentParser(description="Estimate AI-likeness of text from various file types.")
    parser.add_argument("filepath", help="Path to input file (.txt, .docx, .pdf, .pptx)")
    args = parser.parse_args()

    text = extract_text_from_file(args.filepath)
    analyze_text(text)

if __name__ == "__main__":
    main()
